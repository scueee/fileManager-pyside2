import docx
import xlrd # version:1.2.0
from pptx import Presentation
import pdfplumber

from PySide2.QtCore import Signal, QObject, Slot, QDirIterator, QFile, QIODevice, QDir
import re
import win32com.client  # pypiwin32
# import comtypes.client

class Worker(QObject):
    table_view_signal = Signal(int, str, str)
    file_preview_signal = Signal(str)
    progress_bar_signal = Signal(int, int)

    @Slot(str)
    def fileSearch(self, isNameSearch:bool, path:str, text:str):
        p_c = re.compile(text)
        dirIt = QDirIterator(path, QDirIterator.Subdirectories)  # 获取目录信息
        if isNameSearch:
            i = 0
            while dirIt.hasNext():
                dirIt.next()
                if re.search(p_c, dirIt.fileName()):
                    self.table_view_signal.emit(i, dirIt.fileName(), dirIt.filePath())
                    i += 1
        else:
            i = 0
            while dirIt.hasNext():
                dirIt.next()
                path = dirIt.filePath()
                self.text = ""
                self.word = win32com.client.DispatchEx('Word.Application')
                self.pptapp = win32com.client.DispatchEx('PowerPoint.Application')
                self.word.Visible = 0  # 后台运行,不显示
                self.word.DisplayAlerts = 0  # 不警告
                if re.search(".*(pdf|PDF){1}$", path):
                    with pdfplumber.open(path) as pdf:
                        for page in pdf.pages:
                            self.text += page.extract_text()
                elif re.search(".*(docx|DOCX){1}$", path):
                    data = docx.Document(path).paragraphs
                    for p in data:
                        self.text += p.text
                elif re.search(".*(doc|DOC){1}$", path):
                    doc = self.word.Documents.Open(path)
                    for item in doc.paragraphs:
                        self.text += item.Range.Text
                elif re.search(".*(xls|xlsx|XLS|XLSX){1}$", path):
                    data = xlrd.open_workbook(path)
                    sheets = data.sheets()
                    for sheet in sheets:
                        n = sheet.nrows
                        for i in range(n):
                            self.text += "".join(map(lambda x: str(x), sheet.row_values(i)))
                elif re.search(".*(pptx|PPTX){1}$", path):
                    ppt = Presentation(path)
                    for slide in ppt.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                self.text += shape.text_frame.text
                elif re.search(".*(ppt|PPT){1}$", path):
                    ppt = self.pptapp.Presentations.Open(path)
                    page_count = ppt.Slides.Count
                    for i in range(1,page_count+1):
                        shape_count = ppt.Slides(i).Shapes.Count
                        for j in range(1,shape_count+1):
                            if ppt.Slides(i).Shapes(j).HasTextFrame:
                                self.text += ppt.Slides(i).Shapes(j).TextFrame.TextRange.Text
                else:
                    f = QFile(path)
                    f.open(QIODevice.ReadOnly)
                    text = f.readAll()
                    self.text = str(text, 'utf8')
                    f.close()
                self.word.Quit()
                self.pptapp.Quit()
                print("get text")
                if re.search(p_c, self.text):
                    self.table_view_signal.emit(i, dirIt.fileName(), dirIt.filePath())
                    i += 1


    @Slot(str)
    def postText(self, path: str):
        if re.search(".*(pdf|PDF){1}$", path):
            self.postPdfText(path)
        elif re.search(".*(docx|DOCX){1}$", path):
            self.postDocxText(path)
        elif re.search(".*(doc|DOC){1}$", path):
            self.postDocText(path)
        elif re.search(".*(xls|xlsx|XLS|XLSX){1}$", path):
            self.postExcelText(path)
        elif re.search(".*(pptx|PPTX){1}$", path):
            self.postPptxText(path)
        elif re.search(".*(ppt|PPT){1}$", path):
            self.postPptText(path)
        else:
            self.postText(path)

    def postPdfText(self, path: str):
        with pdfplumber.open(path) as pdf:
            progMax = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text()
                progValue = page.page_number
                self.file_preview_signal.emit(text)
                self.progress_bar_signal.emit(progMax, progValue)

    def postDocxText(self, path: str):
        data = docx.Document(path).paragraphs
        progMax = len(data)-1
        for idx, p in enumerate(data):
            self.progress_bar_signal.emit(progMax, idx)
            self.file_preview_signal.emit(f"{p.text}\n")

    def postDocText(self, path: str):
        self.app = win32com.client.DispatchEx('Word.Application')  # 打开独立进程word应用程序
        # self.app = comtypes.client.CreateObject('Word.Application')
        self.app.Visible = 0  # 后台运行,不显示
        self.app.DisplayAlerts = 0  # 不警告
        # doc = self.app.Documents.Open(FileName=path, Encoding='gbk')
        doc = self.app.Documents.Open(path)
        progMax = len(doc.paragraphs) - 1
        for idx, item in enumerate(doc.paragraphs):
            self.progress_bar_signal.emit(progMax, idx)
            self.file_preview_signal.emit(item.Range.Text)
        # for t in doc.Tables:
        #     for row in t.Rows:
        #         for cell in row.Cells:
        #             print(cell.Range.Text)
        self.app.Documents.Close(SaveChanges=0)
        self.app.Quit()

    def postExcelText(self, path: str):
        data = xlrd.open_workbook(path)
        sheets = data.sheets()
        for sheet in sheets:
            n = sheet.nrows
            for i in range(n):
                text = " | ".join(map(lambda x: str(x), sheet.row_values(i)))
                self.file_preview_signal.emit(f"{text}\n")
                self.progress_bar_signal.emit(n, i)
        self.progress_bar_signal.emit(100, 100)

    def postPptxText(self, path: str):
        ppt = Presentation(path)
        progMax = len(ppt.slides)-1
        for idx, slide in enumerate(ppt.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.file_preview_signal.emit(shape.text_frame.text)
            self.progress_bar_signal.emit(progMax, idx)

    def postPptText(self, path: str):
        self.app = win32com.client.DispatchEx('PowerPoint.Application')  # 打开独立进程应用
        self.app.Visible = 0  # 后台运行,不显示
        self.app.DisplayAlerts = 0  # 不警告
        ppt = self.app.Presentations.Open(path)
        page_count = ppt.Slides.Count
        for i in range(1,page_count+1):
            shape_count = ppt.Slides(i).Shapes.Count
            text = ""
            for j in range(1,shape_count+1):
                if ppt.Slides(i).Shapes(j).HasTextFrame:
                    text += ppt.Slides(i).Shapes(j).TextFrame.TextRange.Text
            self.file_preview_signal.emit(text)
            self.progress_bar_signal.emit(page_count, i)
        self.app.Presentations.Close()
        self.app.Quit()

    def postText(self, path:str):
        f = QFile(path)
        f.open(QIODevice.ReadOnly)
        text = f.readAll()
        self.file_preview_signal.emit(str(text,'utf8'))
        f.close()

if __name__ == "__main__":
    worker = Worker()
    worker.postWordText(r"C:\Users\14646\Downloads\b.doc")